import os
import re
import subprocess
from datetime import datetime
from pathlib import Path
from urllib.parse import parse_qs, quote_plus, urlparse

from flask import Flask, render_template, request
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt
import requests
from bs4 import BeautifulSoup

app = Flask(__name__)

# Local mode: each user chooses where the PPT is saved on their own computer.
DEFAULT_OUTPUT_DIR = Path(
    os.getenv("DEFAULT_OUTPUT_DIR", str(Path.home() / "Downloads" / "Generated-PPT-Decks"))
).resolve()
TEMPLATE_PATH = Path(
    os.getenv("PPT_TEMPLATE_PATH", "../../docs/brand/EPAM_PresalesTemplate.potx")
).resolve()
LOGO_PATH = Path(
    os.getenv("LOGO_PATH", "../../docs/brand/EPAM_LOGO_Black.png")
).resolve()

# EPAM brand palette.
BRAND = {
    "cobalt": RGBColor(0, 71, 255),
    "aqua": RGBColor(0, 120, 194),
    "iris": RGBColor(132, 83, 210),
    "snow": RGBColor(251, 250, 250),
    "night": RGBColor(6, 6, 6),
}


def extract_insights(research: list[dict], max_points: int = 5) -> list[str]:
    insights: list[str] = []
    seen: set[str] = set()

    for item in research:
        raw = (item.get("summary") or item.get("snippet") or "").strip()
        if not raw:
            continue

        for sentence in re.split(r"(?<=[.!?])\s+", raw):
            text = re.sub(r"\s+", " ", sentence).strip(" -\t\n\r")
            if len(text) < 70 or len(text) > 180:
                continue

            norm = re.sub(r"[^a-z0-9]+", "", text.lower())
            if not norm or norm in seen:
                continue

            seen.add(norm)
            insights.append(text)
            if len(insights) >= max_points:
                return insights

    return insights


def convert_potx_to_pptx(potx_path: Path) -> Path | None:
    """Convert .potx to .pptx using local PowerPoint COM automation on Windows."""
    if potx_path.suffix.lower() != ".potx" or not potx_path.exists():
        return None

    cache_dir = Path("./template-cache").resolve()
    cache_dir.mkdir(parents=True, exist_ok=True)
    out_pptx = cache_dir / f"{potx_path.stem}.runtime.pptx"

    # Reuse the converted file when it is up-to-date.
    if out_pptx.exists() and out_pptx.stat().st_mtime >= potx_path.stat().st_mtime:
        return out_pptx

    src = str(potx_path).replace("'", "''")
    dst = str(out_pptx).replace("'", "''")
    ps_script = (
        "$ErrorActionPreference='Stop';"
        f"$src='{src}';$dst='{dst}';"
        "$ppt=New-Object -ComObject PowerPoint.Application;"
        "$pres=$ppt.Presentations.Open($src,$false,$false,$false);"
        "$pres.SaveAs($dst,24);"
        "$pres.Close();"
        "$ppt.Quit();"
    )

    try:
        result = subprocess.run(
            [
                "powershell",
                "-NoProfile",
                "-NonInteractive",
                "-ExecutionPolicy",
                "Bypass",
                "-Command",
                ps_script,
            ],
            capture_output=True,
            text=True,
            timeout=90,
            check=False,
        )
        if result.returncode == 0 and out_pptx.exists():
            return out_pptx
    except Exception:
        return None

    return None


def resolve_duckduckgo_link(raw_url: str) -> str:
    parsed = urlparse(raw_url)
    if parsed.netloc.endswith("duckduckgo.com") and parsed.path.startswith("/l/"):
        q = parse_qs(parsed.query)
        if "uddg" in q and q["uddg"]:
            return q["uddg"][0]
    return raw_url


def safe_domain(url: str) -> str:
    try:
        return urlparse(url).netloc.lower()
    except Exception:
        return ""


def web_search(query: str, max_results: int = 6) -> list[dict]:
    search_url = f"https://duckduckgo.com/html/?q={quote_plus(query)}"
    headers = {"User-Agent": "Mozilla/5.0"}
    resp = requests.get(search_url, headers=headers, timeout=20)
    resp.raise_for_status()

    soup = BeautifulSoup(resp.text, "html.parser")
    out: list[dict] = []
    seen_urls: set[str] = set()

    for result in soup.select("div.result"):
        link_el = result.select_one("a.result__a")
        if not link_el:
            continue

        raw_url = link_el.get("href", "").strip()
        url = resolve_duckduckgo_link(raw_url)
        if not url or url in seen_urls:
            continue

        snippet_el = result.select_one(".result__snippet")
        snippet = snippet_el.get_text(" ", strip=True) if snippet_el else ""

        out.append(
            {
                "title": link_el.get_text(" ", strip=True)[:180],
                "url": url,
                "snippet": snippet[:320],
            }
        )
        seen_urls.add(url)

        if len(out) >= max_results:
            break

    return out


def fetch_page_summary(url: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0"}
    try:
        resp = requests.get(url, headers=headers, timeout=20)
        resp.raise_for_status()
    except Exception:
        return ""

    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.extract()

    for p in soup.find_all("p"):
        text = re.sub(r"\s+", " ", p.get_text(" ", strip=True)).strip()
        if len(text) >= 90:
            return text[:320]

    return ""


def gather_research(topic: str, industry: str, audience: str) -> list[dict]:
    queries = [
        f"{topic} {industry} trends 2025 2026",
        f"{topic} {industry} market statistics",
        f"{topic} best practices for {audience}",
    ]

    pooled: list[dict] = []
    for q in queries:
        try:
            pooled.extend(web_search(q, max_results=4))
        except Exception:
            continue

    # Keep diverse domains and avoid duplicates.
    selected: list[dict] = []
    used_domains: set[str] = set()
    used_urls: set[str] = set()

    for item in pooled:
        url = item.get("url", "")
        domain = safe_domain(url)
        if not url or url in used_urls:
            continue

        if domain in used_domains and len(selected) >= 4:
            continue

        summary = fetch_page_summary(url)
        if summary:
            item["summary"] = summary

        selected.append(item)
        used_urls.add(url)
        if domain:
            used_domains.add(domain)

        if len(selected) >= 6:
            break

    if selected:
        return selected

    # Fallback: use Wikipedia API so we still provide online, citeable sources.
    wiki_candidates = [
        topic,
        f"{topic} in {industry}",
        industry,
    ]

    wiki_results: list[dict] = []
    headers = {"User-Agent": "Mozilla/5.0"}
    for candidate in wiki_candidates:
        title = candidate.strip().replace(" ", "_")
        if not title:
            continue
        url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{quote_plus(title)}"
        try:
            resp = requests.get(url, headers=headers, timeout=20)
            if resp.status_code != 200:
                continue
            payload = resp.json()
            if payload.get("type") == "https://mediawiki.org/wiki/HyperSwitch/errors/not_found":
                continue

            page_url = payload.get("content_urls", {}).get("desktop", {}).get("page", "")
            summary = (payload.get("extract") or "").strip()
            page_title = (payload.get("title") or candidate).strip()
            if page_url and summary:
                wiki_results.append(
                    {
                        "title": page_title,
                        "url": page_url,
                        "snippet": summary[:320],
                        "summary": summary[:320],
                    }
                )
        except Exception:
            continue

    if len(wiki_results) < 2:
        # Broader fallback: Wikipedia OpenSearch returns related page URLs.
        for candidate in wiki_candidates:
            try:
                api_url = (
                    "https://en.wikipedia.org/w/api.php?action=opensearch"
                    f"&search={quote_plus(candidate)}&limit=8&namespace=0&format=json"
                )
                resp = requests.get(api_url, headers=headers, timeout=20)
                if resp.status_code != 200:
                    continue

                payload = resp.json()
                titles = payload[1] if len(payload) > 1 else []
                urls = payload[3] if len(payload) > 3 else []
                for title, page_url in zip(titles, urls):
                    if not page_url:
                        continue
                    if any(existing.get("url") == page_url for existing in wiki_results):
                        continue

                    wiki_results.append(
                        {
                            "title": str(title),
                            "url": str(page_url),
                            "snippet": f"Reference page for {title}",
                            "summary": f"Reference page for {title}",
                        }
                    )
                    if len(wiki_results) >= 4:
                        break
                if len(wiki_results) >= 4:
                    break
            except Exception:
                continue

    return wiki_results


def choose_template_path() -> tuple[Path | None, str]:
    # Prefer usable .pptx templates; if only .potx exists, convert it.
    candidates = [
        TEMPLATE_PATH,
        Path("../../docs/brand/EPAM_PresalesTemplate.pptx").resolve(),
        Path("../../docs/brand/EPAM_PresalesTemplate.potx").resolve(),
    ]

    for candidate in candidates:
        if not candidate.exists():
            continue

        suffix = candidate.suffix.lower()
        if suffix == ".pptx":
            return candidate, f"Using template: {candidate.name}"
        if suffix == ".potx":
            converted = convert_potx_to_pptx(candidate)
            if converted:
                return converted, f"Converted template from {candidate.name}"

    return None, "No compatible template found. Used branded fallback layout."


def safe_slug(value: str) -> str:
    value = value.strip().lower()
    value = re.sub(r"[^a-z0-9]+", "-", value)
    return value.strip("-") or "capability"


def style_title(slide) -> None:
    if slide.shapes.title is None:
        return
    tf = slide.shapes.title.text_frame
    tf.word_wrap = True
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(28)
        run.font.color.rgb = BRAND["snow"]
    try:
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    except Exception:
        pass


def add_title_bar(slide, prs: Presentation) -> None:
    bar_height = Inches(0.7)
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        0,
        0,
        prs.slide_width,
        bar_height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND["cobalt"]
    shape.line.fill.background()


def add_logo(slide, prs: Presentation) -> None:
    if not LOGO_PATH.exists():
        return
    width = Inches(1.2)
    left = prs.slide_width - width - Inches(0.2)
    top = Inches(0.08)
    slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width)


def add_sources_to_notes(slide, sources: list[dict]) -> None:
    try:
        tf = slide.notes_slide.notes_text_frame
        tf.clear()
        tf.text = "Sources"
        for src in sources:
            p = tf.add_paragraph()
            p.text = f"- {src.get('title', 'Source')}: {src.get('url', '')}"
    except Exception:
        # Notes are best effort.
        pass


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], sources: list[dict] | None = None):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    add_title_bar(slide, prs)

    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        style_title(slide)
        add_logo(slide, prs)

    body = None
    for ph in slide.placeholders:
        if slide.shapes.title is not None and ph == slide.shapes.title:
            continue
        if hasattr(ph, "text_frame"):
            body = ph.text_frame
            break

    if body is None:
        # Fallback for templates that do not expose a standard body placeholder.
        tx = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.2),
            prs.slide_width - Inches(1.6),
            prs.slide_height - Inches(1.8),
        )
        body = tx.text_frame

    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        try:
            p.runs[0].font.size = Pt(14)
            p.runs[0].font.color.rgb = BRAND["night"]
        except Exception:
            pass

    if sources:
        add_sources_to_notes(slide, sources)

    return slide


def generate_deck(
    topic: str,
    audience: str,
    industry: str,
    extra_prompt: str,
    output_dir: Path,
) -> tuple[Path, list[dict], str]:
    template_path, template_status = choose_template_path()
    if template_path:
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()

    research = gather_research(topic, industry, audience)
    insights = extract_insights(research, max_points=6)

    market_bullets = [point[:150] for point in insights[:3]]

    if not market_bullets:
        market_bullets = [
            f"{industry} leaders are increasing focus on {topic} outcomes.",
            "Executives are prioritizing measurable value, speed, and governance.",
            "Teams need scalable execution models with clear ROI visibility.",
        ]

    capability_bullets = [
        f"Advisory-to-delivery model for {topic} programs in {industry}.",
        f"Value tracking aligned to {audience} priorities and business KPIs.",
        "Reusable accelerators to reduce time-to-value and delivery risk.",
    ]

    if insights:
        capability_bullets.append(insights[0][:150])

    proof_points_bullets = [
        "Reduced cycle time through automation and operating model redesign.",
        "Improved quality through governance controls and data standards.",
        "Enabled faster decisions with trustworthy, role-specific reporting.",
    ]

    if len(insights) > 1:
        proof_points_bullets[2] = insights[1][:150]

    today = datetime.now().strftime("%Y-%m-%d")
    file_name = f"{safe_slug(topic)}-deck-{today}.pptx"

    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / file_name

    # Quality deck baseline with sources and executive narrative.
    add_bullet_slide(prs, f"{topic} Capability", [
        f"Audience: {audience}",
        f"Industry: {industry}",
        "Prepared with current market signals, sources, and branded storytelling",
    ])
    add_bullet_slide(prs, "Agenda", [
        "Market context",
        "Our capability",
        "How we work",
        "Proof points",
        "Sources and next steps",
    ])
    add_bullet_slide(prs, "Market Context", market_bullets, sources=research[:3])
    add_bullet_slide(prs, "Our Capability", capability_bullets, sources=research[:2])
    add_bullet_slide(prs, "How We Work", [
        "Diagnose current state, constraints, and decision bottlenecks.",
        "Define target operating model, architecture, and KPI framework.",
        "Pilot in a focused scope to prove value and derisk change.",
        "Scale in waves with governance, enablement, and adoption plans.",
        "Continuously optimize outcomes using KPI and benefit tracking.",
    ])
    add_bullet_slide(prs, "Proof Points", proof_points_bullets, sources=research[:3])
    add_bullet_slide(prs, "Why Us", [
        f"Cross-functional experts across strategy, delivery, and change in {industry}.",
        "Battle-tested accelerators and templates to shorten execution timelines.",
        "Governed delivery model focused on measurable business outcomes.",
    ])
    add_bullet_slide(prs, "Next Steps", [
        "Run discovery workshop",
        "Agree 90-day roadmap",
        "Launch first execution sprint",
    ])
    source_lines = [
        f"{item.get('title', 'Source')} - {item.get('url', '')}" for item in research[:5]
    ]
    add_bullet_slide(
        prs,
        "Sources",
        source_lines if source_lines else ["No online sources available at generation time."],
        sources=research,
    )
    add_bullet_slide(prs, "Thank You", [
        "Contact the consulting team",
        "Deck created with brand styling and source-backed context",
        f"Notes: {extra_prompt[:120] or 'N/A'}",
    ])

    prs.save(str(out_path))
    return out_path, research, template_status


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html", default_output_dir=str(DEFAULT_OUTPUT_DIR))


@app.route("/generate", methods=["POST"])
def generate():
    topic = request.form.get("topic", "").strip()
    audience = request.form.get("audience", "").strip()
    industry = request.form.get("industry", "").strip()
    extra_prompt = request.form.get("extra_prompt", "").strip()
    output_folder = request.form.get("output_folder", "").strip()

    chosen_output_dir = Path(output_folder).expanduser().resolve() if output_folder else DEFAULT_OUTPUT_DIR

    if not topic or not audience or not industry or not output_folder:
        return render_template(
            "index.html",
            error="Topic, audience, industry, and output folder are required.",
            form=request.form,
            default_output_dir=str(DEFAULT_OUTPUT_DIR),
        )

    out_path, research, template_status = generate_deck(
        topic,
        audience,
        industry,
        extra_prompt,
        chosen_output_dir,
    )

    return render_template(
        "index.html",
        success=True,
        file_name=out_path.name,
        local_file_path=str(out_path),
        saved_folder=str(chosen_output_dir),
        default_output_dir=str(DEFAULT_OUTPUT_DIR),
        sources=research,
        template_status=template_status,
        logo_status=("Logo applied" if LOGO_PATH.exists() else "Logo file not found"),
        quality_status=(
            f"Quality gate passed: {len(research)} sources captured"
            if len(research) >= 2
            else "Quality gate warning: fewer than 2 sources captured"
        ),
    )


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=8000, debug=False)
